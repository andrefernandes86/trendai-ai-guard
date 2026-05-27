import io
import json
import logging
import re

logger = logging.getLogger(__name__)


def extract_text(file_bytes: bytes, ext: str) -> str:
    """Return the text content of a file given its raw bytes and extension."""
    dispatch = {
        ".txt": _plain,
        ".csv": _plain,
        ".md": _plain,
        ".xml": _plain,
        ".yaml": _plain,
        ".yml": _plain,
        ".html": _plain,
        ".htm": _plain,
        ".json": _json,
        ".pdf": _pdf,
        ".docx": _docx,
        ".doc": _docx,
        ".xlsx": _xlsx,
        ".xls": _xlsx,
        ".pptx": _pptx,
        ".ppt": _pptx,
        ".rtf": _rtf,
    }
    handler = dispatch.get(ext, _plain)
    try:
        return handler(file_bytes)
    except Exception:
        logger.exception("Text extraction failed for extension %s", ext)
        return ""


def _plain(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


def _json(data: bytes) -> str:
    try:
        return json.dumps(json.loads(data.decode("utf-8")), indent=2)
    except Exception:
        return _plain(data)


def _pdf(data: bytes) -> str:
    try:
        import pypdf  # noqa: PLC0415
        reader = pypdf.PdfReader(io.BytesIO(data))
        parts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(parts)
    except ImportError:
        logger.error("pypdf is not installed")
        return ""


def _docx(data: bytes) -> str:
    try:
        from docx import Document  # noqa: PLC0415
        doc = Document(io.BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except ImportError:
        logger.error("python-docx is not installed")
        return ""


def _xlsx(data: bytes) -> str:
    try:
        import openpyxl  # noqa: PLC0415
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                line = " ".join(str(c) for c in row if c is not None)
                if line.strip():
                    rows.append(line)
        return "\n".join(rows)
    except ImportError:
        logger.error("openpyxl is not installed")
        return ""


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # noqa: PLC0415
        prs = Presentation(io.BytesIO(data))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except ImportError:
        logger.error("python-pptx is not installed")
        return ""


def _rtf(data: bytes) -> str:
    try:
        from striprtf.striprtf import rtf_to_text  # noqa: PLC0415
        return rtf_to_text(data.decode("utf-8", errors="replace"))
    except ImportError:
        # Best-effort fallback: strip RTF control words
        text = data.decode("utf-8", errors="replace")
        text = re.sub(r"\{[^{}]*\}", "", text)
        text = re.sub(r"\\[a-z]+\d* ?", "", text)
        return text.strip()
